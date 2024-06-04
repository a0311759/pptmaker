
import uuid
from flask import Flask, render_template, request, jsonify
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
import os
from datetime import datetime

app = Flask(__name__)

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/download', methods=['POST'])
def download():
    data = request.json
    img_option = data.get('img_option', 'a')
    print("Selected image option:", img_option)

    modified_slides = []
    for slide in data.get('slidesData', []):
        heading = slide['heading'].replace('\xa0', ' ')
        paragraph = slide['paragraph'].replace('\xa0', ' ')
        if paragraph != "Editable Paragraph 1: Lorem ipsum dolor sit amet, consectetur adipiscing elit. Sed do eiusmod tempor incididunt ut labore et dolore magna aliqua. Ut enim ad minim veniam, quis nostrud exercitation ullamco laboris nisi ut aliquip ex ea commodo consequat. Duis aute irure dolor in reprehenderit in voluptate velit esse cillum dolore eu fugiat nulla pariatur. Excepteur sint occaecat cupidatat non proident, sunt in culpa qui officia deserunt mollit anim id est laborum.":
            modified_slides.append({'heading': heading, 'paragraph': paragraph})
    
    directory = r'E:\projects\presentation\quick-slides\static'
    img_paths = {
        'a': "E:\\projects\\presentation\\quick-slides\\images\\bubble.jpg",
        'b': "E:\\projects\\presentation\\quick-slides\\images\\pink.jpg",
        'c': "E:\\projects\\presentation\\quick-slides\\images\\fold-paper.jpg",
        'd': "E:\\projects\\presentation\\quick-slides\\images\\butter-fly.jpg"
    }
    img_path = img_paths.get(img_option, "E:\\projects\\presentation\\bubble.jpg")

    prs = Presentation()
    for i, slide_data in enumerate(modified_slides):
        if i == 0:
            slide_layout = prs.slide_layouts[0]
        else:
            slide_layout = prs.slide_layouts[1]
            
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
        
        title = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        content = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
        title.text = slide_data['heading']
        words = slide_data['paragraph'].split()
        lines = [' '.join(words[i:i+8]) for i in range(0, len(words), 8)]
        content.text = '\n'.join(lines)
        
        if i == 0:
            for paragraph in title.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
            for paragraph in content.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
        else:
            for paragraph in title.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
        for paragraph in content.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
        
        title.text_frame.paragraphs[0].font.size = Pt(56 if i == 0 else 32)
        title.text_frame.paragraphs[0].font.bold = True
        content.text_frame.paragraphs[0].font.size = Pt(36 if i == 0 else 20)
        
        content.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        
    unique_id = uuid.uuid4().hex[:6]
    current_time = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = os.path.join(directory, f'presentation_{current_time}_{unique_id}.pptx')
    
    prs.save(filename)
    print(f"Presentation saved as: {filename}")
    
    return jsonify({'filename': filename})

if __name__ == '__main__':
    app.run(debug=True)
